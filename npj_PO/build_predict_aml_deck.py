# -*- coding: utf-8 -*-
"""
PREDICT-AML storytelling deck builder.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from PIL import Image
import copy
import os

HERE = os.path.dirname(os.path.abspath(__file__))
MEDIA = os.path.join(HERE, "pptx_assets")
OUT = os.path.join(HERE, "PREDICT-AML_Presentation.pptx")

# ---------------------------------------------------------------- palette --
NAVY      = RGBColor(0x0D, 0x1B, 0x38)
NAVY2     = RGBColor(0x16, 0x28, 0x4D)
INK       = RGBColor(0x15, 0x22, 0x38)
BG_LIGHT  = RGBColor(0xFA, 0xFA, 0xF8)
CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x6B, 0x76, 0x8A)
GRAY_LT   = RGBColor(0xD8, 0xDC, 0xE4)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GOLD_TXT  = RGBColor(0xF2, 0xB2, 0x33)

CORAL   = RGBColor(0xE4, 0x57, 0x2E)   # Act I   - The Problem
TEAL    = RGBColor(0x21, 0x9E, 0x8B)   # Act II  - The Framework
BLUE    = RGBColor(0x2E, 0x62, 0x9E)   # Act III - Putting it to the test
PURPLE  = RGBColor(0x6A, 0x4C, 0x93)   # Act IV  - The Biology Inside
GOLD    = RGBColor(0xD1, 0x8A, 0x2E)   # Act V   - Proof & Payoff
GREEN   = RGBColor(0x2E, 0x93, 0x3C)

ACTS = {
    1: ("THE PROBLEM", CORAL),
    2: ("THE FRAMEWORK", TEAL),
    3: ("PUTTING IT TO THE TEST", BLUE),
    4: ("THE BIOLOGY INSIDE", PURPLE),
    5: ("PROOF & PAYOFF", GOLD),
}

FONT = "Calibri"
FONT_H = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

page_counter = {"n": 0}

# ------------------------------------------------------------- utilities --
def add_slide():
    s = prs.slides.add_slide(BLANK)
    page_counter["n"] += 1
    return s

def rect(slide, x, y, w, h, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh

def rrect(slide, x, y, w, h, color, radius=0.08, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh

def oval(slide, x, y, w, h, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    return sh

def textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False, italic=False,
            align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
            wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb

def rich_textbox(slide, x, y, w, h, runs_per_line, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.0, wrap=True):
    """runs_per_line: list of list-of-dict(text,size,color,bold,italic,font)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, runs in enumerate(runs_per_line):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for rd in runs:
            r = p.add_run()
            r.text = rd["text"]
            r.font.size = Pt(rd.get("size", 18))
            r.font.bold = rd.get("bold", False)
            r.font.italic = rd.get("italic", False)
            r.font.color.rgb = rd.get("color", INK)
            r.font.name = rd.get("font", FONT)
    return tb

def picture_fit(slide, path, x, y, w, h):
    """Place image scaled to fit inside box, preserving aspect ratio, centered."""
    im = Image.open(path)
    iw, ih = im.size
    box_ratio = w / h
    im_ratio = iw / ih
    if im_ratio > box_ratio:
        nw = w
        nh = int(w / im_ratio)
    else:
        nh = h
        nw = int(h * im_ratio)
    nx = x + (w - nw) // 2
    ny = y + (h - nh) // 2
    pic = slide.shapes.add_picture(path, nx, ny, width=nw, height=nh)
    return pic

def kicker_and_bar(slide, act_no, extra=""):
    label, color = ACTS[act_no]
    rect(slide, 0, 0, SLIDE_W, Inches(0.09), color)
    textbox(slide, Inches(0.55), Inches(0.28), Inches(10), Inches(0.35),
            label + (("   ·   " + extra) if extra else ""), size=12.5, color=color,
            bold=True, font=FONT)
    return color

def slide_title(slide, text, color=INK, y=Inches(0.62), size=30, w=Inches(12.2)):
    textbox(slide, Inches(0.55), y, w, Inches(0.9), text, size=size, color=color,
            bold=True, font=FONT_H)

def footer(slide, dark=False):
    n = page_counter["n"]
    c = GRAY_LT if dark else GRAY
    textbox(slide, Inches(0.55), Inches(7.14), Inches(4), Inches(0.3),
            "PREDICT-AML", size=9, color=c, bold=True)
    textbox(slide, Inches(12.4), Inches(7.14), Inches(0.5), Inches(0.3),
            str(n), size=9, color=c, align=PP_ALIGN.RIGHT)

def progress_dots(slide, current_act, dark=False):
    acts_order = [1, 2, 3, 4, 5]
    total = len(acts_order)
    dot_w = Inches(0.32)
    gap = Inches(0.14)
    total_w = total * dot_w + (total - 1) * gap
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(7.12)
    for i, a in enumerate(acts_order):
        _, color = ACTS[a]
        x = start_x + i * (dot_w + gap)
        c = color if a == current_act else (RGBColor(0x33, 0x40, 0x5C) if dark else GRAY_LT)
        h = Inches(0.085)
        rrect(slide, x, y, dot_w, h, c, radius=0.5)

def bg(slide, color=BG_LIGHT):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

def bignum(slide, x, y, w, number, label, color, num_size=54, label_size=13,
           align=PP_ALIGN.LEFT):
    textbox(slide, x, y, w, Inches(0.9), number, size=num_size, color=color, bold=True,
            align=align)
    textbox(slide, x, y + Inches(0.78), w, Inches(0.7), label, size=label_size, color=GRAY,
            align=align, line_spacing=1.05)

def chip(slide, x, y, w, h, text, bg_color, txt_color=WHITE, size=13):
    sh = rrect(slide, x, y, w, h, bg_color, radius=0.5)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = txt_color; r.font.name = FONT
    return sh

def connector_arrow(slide, x1, y1, x2, y2, color=GRAY, width=1.5):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    line = cn.line._get_or_add_ln()
    tail = line.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    line.append(tail)
    return cn

def card(slide, x, y, w, h, fill=CARD_BG, line=GRAY_LT):
    return rrect(slide, x, y, w, h, fill, radius=0.06, line=line)

print("helpers module executed OK")

# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, SLIDE_W, Inches(0.09), GOLD_TXT)

textbox(s, Inches(0.7), Inches(1.15), Inches(8.5), Inches(0.4),
        "A  PRECISION-ONCOLOGY  AI  FRAMEWORK", size=14, color=GOLD_TXT, bold=True)
textbox(s, Inches(0.65), Inches(1.55), Inches(9.5), Inches(1.5),
        "PREDICT-AML", size=64, color=WHITE, bold=True)
textbox(s, Inches(0.7), Inches(2.75), Inches(8.6), Inches(1.0),
        "Personalized REsponse via Drug Interaction with Cellular Traits\nfor Acute Myeloid Leukemia",
        size=19, color=RGBColor(0xC7, 0xD1, 0xE6), line_spacing=1.15)

rect(s, Inches(0.72), Inches(3.95), Inches(0.55), Inches(0.04), GOLD_TXT)
textbox(s, Inches(0.7), Inches(4.15), Inches(9.2), Inches(0.8),
        "Teaching a model to matchmake 805 leukemia patients with 165 drugs —\n"
        "and to be honest about when it actually works.",
        size=15.5, color=GOLD_TXT, italic=True, line_spacing=1.2)

textbox(s, Inches(0.7), Inches(6.55), Inches(9), Inches(0.5),
        "Mohammed Al-Ani   ·   Siddhi P. Jani   ·   Halima Bensmail   ·   Raghvendra Mall",
        size=13.5, color=WHITE, bold=True)
textbox(s, Inches(0.7), Inches(6.9), Inches(9), Inches(0.4),
        "Qatar Computing Research Institute, HBKU", size=11.5, color=GRAY_LT)

# --- decorative "matchmaking" graphic, right side ---
cx = Inches(11.05)
pill = rrect(s, Inches(10.35), Inches(1.55), Inches(1.4), Inches(0.62), WHITE, radius=0.5)
rrect(s, Inches(10.35)+Inches(0.7), Inches(1.55), Inches(0.7), Inches(0.62), GOLD_TXT, radius=0.0)
# re-round only the right half by overlaying rounded rect clipped visually is complex; keep simple pill + divider line
ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.35)+Inches(0.7), Inches(1.55),
                             Inches(10.35)+Inches(0.7), Inches(2.17))
ln.line.color.rgb = NAVY
ln.line.width = Pt(1.5)

colors_dots = [CORAL, GOLD_TXT, TEAL, PURPLE, RGBColor(0x4C, 0x8B, 0xC9)]
sizes = [0.5, 0.62, 0.4, 0.7, 0.46]
yoff = [0.15, 0.02, 0.28, -0.05, 0.2]
xstart = 10.15
for i, (c, sz, yo) in enumerate(zip(colors_dots, sizes, yoff)):
    x = Inches(xstart + i * 0.58)
    y = Inches(2.85 + yo)
    oval(s, x, y, Inches(sz), Inches(sz), c)
    connector_arrow(s, Inches(10.35)+Inches(0.35), Inches(2.17), x + Inches(sz/2), y,
                     color=RGBColor(0x33, 0x40, 0x5C), width=1)

textbox(s, Inches(9.95), Inches(3.7), Inches(2.9), Inches(0.6),
        "one drug → many biological fates", size=10.5, color=GRAY_LT, italic=True,
        align=PP_ALIGN.CENTER)

footer(s, dark=True)

# ============================================================================
# SLIDE 2 — THE CLINICAL PROBLEM (Act I)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 1, "Why precision dosing matters")
slide_title(s, "Acute Myeloid Leukemia doesn't read the textbook")

textbox(s, Inches(0.55), Inches(1.55), Inches(11.5), Inches(0.6),
        "AML is common enough to matter, and lethal enough that guessing the wrong drug costs time patients don't have.",
        size=15, color=GRAY, line_spacing=1.15)

stats = [
    ("22,010", "new AML cases per year\nin the U.S. (SEER, 2025)", CORAL),
    ("11,090", "deaths per year —\nnearly 1 for every 2 diagnoses", CORAL),
    ("60+", "patient subtypes by genomic\nand transcriptomic profile", CORAL),
    ("165", "candidate drugs profiled\nin this study alone", TEAL),
]
card_w = Inches(2.85); gap = Inches(0.22); startx = Inches(0.55); y0 = Inches(2.55)
for i, (num, lab, col) in enumerate(stats):
    x = startx + i * (card_w + gap)
    card(s, x, y0, card_w, Inches(2.55))
    rect(s, x, y0, card_w, Inches(0.09), col)
    bignum(s, x + Inches(0.25), y0 + Inches(0.5), card_w - Inches(0.5), num, lab, col,
           num_size=42, label_size=12.5)

textbox(s, Inches(0.55), Inches(5.5), Inches(11.6), Inches(1.2),
        "The hard truth: two patients with the same AML diagnosis can have completely different tumor biology —"
        " and therefore need completely different drugs. Clinicians currently choose largely by trial, error, and experience.",
        size=15, color=INK, line_spacing=1.25)
progress_dots(s, 1)
footer(s)

print("slides 1-2 built")

def mini_person(slide, x, y, size, color):
    head = int(size * 0.34)
    oval(slide, x + (size - head)//2, y, head, head, color)
    body_w = int(size * 0.8); body_h = int(size * 0.62)
    rrect(slide, x + (size - body_w)//2, y + head - Emu(0), body_w, body_h, color, radius=0.35)

# ============================================================================
# SLIDE 3 — ONE DRUG, MANY FATES  (Act I)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 1, "Same diagnosis, different biology")
slide_title(s, "One drug. Same dose. Wildly different fates.")

textbox(s, Inches(0.55), Inches(1.55), Inches(11.6), Inches(0.55),
        "Standard-of-care assumes AML is one disease. The molecular data says otherwise.",
        size=15, color=GRAY)

# pill icon
pillx, pilly = Inches(6.0), Inches(2.35)
rrect(s, pillx, pilly, Inches(1.3), Inches(0.55), WHITE, radius=0.5, line=GRAY)
rrect(s, pillx + Inches(0.65), pilly, Inches(0.65), Inches(0.55), GOLD, radius=0.0)
connector_arrow(s, pillx + Inches(0.65), pilly, pillx + Inches(0.65), pilly + Inches(0.55), color=GRAY, width=0)
textbox(s, pillx - Inches(0.5), pilly - Inches(0.42), Inches(2.3), Inches(0.35),
        "SAME  DRUG  ·  SAME  DOSE", size=11.5, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

n = 11
row_y = Inches(3.55)
grad_colors = [GREEN, RGBColor(0x4C,0xA9,0x5A), RGBColor(0x8B,0xB8,0x4E), GOLD_TXT,
               RGBColor(0xE3,0x8A,0x35), RGBColor(0xE4,0x57,0x2E), RGBColor(0xE4,0x57,0x2E),
               RGBColor(0xC9,0x3B,0x2E), RGBColor(0x8B,0xB8,0x4E), RGBColor(0xE3,0x8A,0x35), GREEN]
total_w = Inches(11.6)
gap = Emu(int((total_w - n*Inches(0.62)) / (n-1)))
startx = Inches(0.85)
cx_pill = pillx + Inches(0.65)
for i in range(n):
    x = startx + i * (Inches(0.62) + gap)
    y = row_y
    mini_person(s, x, y, Inches(0.62), grad_colors[i])
    connector_arrow(s, cx_pill, pilly + Inches(0.6), x + Inches(0.31), y, color=RGBColor(0xC9,0xCE,0xD8), width=0.75)

textbox(s, Inches(0.85), Inches(4.55), Inches(3.2), Inches(0.4), "◄ SENSITIVE", size=14,
        color=GREEN, bold=True)
textbox(s, Inches(9.25), Inches(4.55), Inches(3.2), Inches(0.4), "RESISTANT ►", size=14,
        color=CORAL, bold=True, align=PP_ALIGN.RIGHT)

card(s, Inches(1.9), Inches(5.35), Inches(9.5), Inches(1.35))
textbox(s, Inches(2.25), Inches(5.55), Inches(8.8), Inches(1.0),
        "Later in this deck: Venetoclax has the widest predicted spread of any drug —\n"
        "mean AUC 149.5, but ±53 SD across BeatAML patients — the single most heterogeneous response we predict.",
        size=14, color=INK, line_spacing=1.25)
progress_dots(s, 1)
footer(s)

# ============================================================================
# SLIDE 4 — BEATAML: THE TREASURE TROVE  (Act I)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 1, "The raw material")
slide_title(s, "BeatAML: a decade of real patient drug-response data")

stats2 = [
    ("805", "patients profiled", CORAL),
    ("942", "tumor specimens", CORAL),
    ("165", "drugs tested ex vivo", TEAL),
    ("~2,000", "multi-omic features\nper patient", TEAL),
]
card_w = Inches(2.85); gap2 = Inches(0.22); startx2 = Inches(0.55); y0 = Inches(1.7)
for i, (num, lab, col) in enumerate(stats2):
    x = startx2 + i * (card_w + gap2)
    card(s, x, y0, card_w, Inches(1.9))
    rect(s, x, y0, card_w, Inches(0.09), col)
    bignum(s, x + Inches(0.25), y0 + Inches(0.35), card_w - Inches(0.5), num, lab, col,
           num_size=34, label_size=12.5)

textbox(s, Inches(0.55), Inches(3.95), Inches(6), Inches(0.4),
        "FOUR  MOLECULAR  LAYERS  PER  PATIENT", size=13, color=INK, bold=True)
layers = [("Gene Expression", "150 most-variable genes\n+ oncogene profiles", CORAL),
          ("Mutations", "somatic variants,\nfrequency & type", TEAL),
          ("Pathway Activity", "GSVA oncogenic\npathway scores", BLUE),
          ("Cell / Module State", "clinical + cell-state +\nmodule enrichment", PURPLE)]
lw = Inches(2.85); gapl = Inches(0.22)
for i, (t, d, c) in enumerate(layers):
    x = startx2 + i * (lw + gapl)
    card(s, x, Inches(4.4), lw, Inches(1.55))
    oval(s, x + Inches(0.18), Inches(4.6), Inches(0.28), Inches(0.28), c)
    textbox(s, x + Inches(0.18), Inches(5.0), lw - Inches(0.36), Inches(0.35), t, size=14.5,
            color=INK, bold=True)
    textbox(s, x + Inches(0.18), Inches(5.4), lw - Inches(0.36), Inches(0.55), d, size=11.5,
            color=GRAY, line_spacing=1.1)

textbox(s, Inches(0.55), Inches(6.25), Inches(11.6), Inches(0.6),
        "Split honestly by time: Waves 1+2 (337 samples) train the model; Waves 3+4 (183 samples),"
        " collected later, are held out as a true, untouched test set.",
        size=13.5, color=GRAY, italic=True, line_spacing=1.2)
progress_dots(s, 1)
footer(s)

# ============================================================================
# SLIDE 5 — LOCKED TREASURE  (Act I)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 1, "What was still missing")
slide_title(s, "Rich data. But nobody had built the key yet.")

gaps = [
    ("Retrospective only", "Findings described the past;\nnothing predicted a new patient's response."),
    ("No stratified benchmark", "Performance was never tested the way\nit would be deployed: on unseen patients."),
    ("Single-cohort scope", "Results lived inside BeatAML,\nnever checked against outside hospitals."),
    ("Drug structure ignored", "Molecular identity of the drug itself\nrarely entered the model at all."),
]
gw = Inches(2.85); gapg = Inches(0.22); startxg = Inches(0.55); y0g = Inches(1.85)
for i, (t, d) in enumerate(gaps):
    x = startxg + i * (gw + gapg)
    card(s, x, y0g, gw, Inches(3.15))
    oval(s, x + Inches(0.22), y0g + Inches(0.22), Inches(0.5), Inches(0.5), CORAL)
    textbox(s, x + Inches(0.22), y0g + Inches(0.28), Inches(0.5), Inches(0.4), "✕", size=20,
            color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x + Inches(0.22), y0g + Inches(0.95), gw - Inches(0.44), Inches(0.65), t, size=15,
            color=INK, bold=True, line_spacing=1.05)
    textbox(s, x + Inches(0.22), y0g + Inches(1.65), gw - Inches(0.44), Inches(1.3), d, size=11.5,
            color=GRAY, line_spacing=1.2)

textbox(s, Inches(0.55), Inches(5.35), Inches(11.6), Inches(0.8),
        "These four gaps define exactly what PREDICT-AML was built to close.",
        size=16, color=TEAL, bold=True, italic=True)
progress_dots(s, 1)
footer(s)

# ============================================================================
# SLIDE 6 — THE FIELD'S BEST ATTEMPTS  (Act I)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 1, "Prior computational attempts")
slide_title(s, "Earlier models got only part of the way there")

chart_data = CategoryChartData()
chart_data.categories = ["ElasticNet\n(Karathanasis '24)", "NetAML\n(per-drug PPI, '25)",
                          "MDREAM\n(drug-specific ens., '23)", "PREDICT-AML\n(this work)"]
chart_data.add_series("Pearson / Spearman r", (0.36, 0.61, 0.68, 0.69))
gx, gy, gw2, gh2 = Inches(0.55), Inches(1.7), Inches(8.6), Inches(4.6)
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, gx, gy, gw2, gh2, chart_data)
chart = gframe.chart
chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
plot.gap_width = 60
series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = GRAY_LT
pts_colors = [CORAL, CORAL, CORAL, GOLD]
for i, pt in enumerate(series.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = pts_colors[i]
cat_ax = chart.category_axis
cat_ax.tick_labels.font.size = Pt(12)
cat_ax.tick_labels.font.color.rgb = INK
val_ax = chart.value_axis
val_ax.minimum_scale = 0
val_ax.maximum_scale = 0.8
val_ax.has_major_gridlines = False
val_ax.tick_labels.font.size = Pt(11)
val_ax.tick_labels.font.color.rgb = GRAY

card(s, Inches(9.4), Inches(1.7), Inches(3.35), Inches(4.6))
textbox(s, Inches(9.65), Inches(1.9), Inches(2.9), Inches(0.4), "WHY  THEY  FELL  SHORT",
        size=12, color=INK, bold=True)
short_txt = [
    "ElasticNet: linear only, 122 drugs, no drug structure",
    "NetAML: 87 separate per-drug models, RNA-seq only, needs ≥35 patients/drug",
    "MDREAM: excludes rare drugs, no uncertainty, no external multi-cohort test",
]
yy = Inches(2.45)
for t in short_txt:
    oval(s, Inches(9.65), yy + Inches(0.06), Inches(0.09), Inches(0.09), CORAL)
    textbox(s, Inches(9.85), yy, Inches(2.75), Inches(0.9), t, size=11.5, color=GRAY, line_spacing=1.15)
    yy += Inches(1.0)

textbox(s, Inches(0.55), Inches(6.45), Inches(8.6), Inches(0.5),
        "Bars use each study's own reported metric/cohort — context, not a leaderboard.",
        size=10.5, color=GRAY, italic=True)
progress_dots(s, 1)
footer(s)

print("slides 1-6 built")

# ============================================================================
# SLIDE 7 — ENTER PREDICT-AML  (Act II, cinematic transition)
# ============================================================================
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, SLIDE_W, Inches(0.09), TEAL)
textbox(s, Inches(0.7), Inches(0.5), Inches(6), Inches(0.4), "ACT II  ·  THE FRAMEWORK",
        size=13, color=TEAL, bold=True)

textbox(s, Inches(0.65), Inches(2.15), Inches(12), Inches(1.3), "Enter PREDICT-AML",
        size=54, color=WHITE, bold=True)
textbox(s, Inches(0.7), Inches(3.35), Inches(10.8), Inches(0.6),
        "One unified model.  Any patient.  Any drug.  Honestly evaluated.",
        size=19, color=GOLD_TXT, italic=True)

eqx = Inches(0.7); eqy = Inches(4.55)
chip(s, eqx, eqy, Inches(3.1), Inches(0.75), "Patient Multi-Omics", CORAL, size=15)
textbox(s, eqx + Inches(3.15), eqy + Inches(0.12), Inches(0.5), Inches(0.55), "+", size=26,
        color=WHITE, bold=True, align=PP_ALIGN.CENTER)
chip(s, eqx + Inches(3.65), eqy, Inches(2.7), Inches(0.75), "Drug Identity", RGBColor(0x4C,0x8B,0xC9), size=15)
textbox(s, eqx + Inches(6.4), eqy + Inches(0.12), Inches(0.5), Inches(0.55), "=", size=26,
        color=WHITE, bold=True, align=PP_ALIGN.CENTER)
chip(s, eqx + Inches(6.95), eqy, Inches(3.4), Inches(0.75), "Predicted Drug Response", GOLD_TXT,
     txt_color=NAVY, size=15)

footer(s, dark=True)

# ============================================================================
# SLIDE 8 — FRAMEWORK OVERVIEW (flowchart)  (Act II)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 2, "The full pipeline, at a glance")
slide_title(s, "From raw multi-omics to an actionable prediction")
picture_fit(s, os.path.join(MEDIA, "Predict-AML_Flowchart-1.png"),
            Inches(0.55), Inches(1.35), Inches(12.25), Inches(5.55))
progress_dots(s, 2)
footer(s)

# ============================================================================
# SLIDE 9 — INGREDIENT 1: THE PATIENT  (Act II)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 2, "Ingredient #1")
slide_title(s, "The patient: four molecular layers, one feature vector")

textbox(s, Inches(0.55), Inches(1.5), Inches(11.6), Inches(0.5),
        "~1,300 engineered features distilled per patient — enough signal, without drowning the model in noise.",
        size=14.5, color=GRAY)

feat_rows = [
    ("Gene Expression", "150 most-variable genes + oncogene panel", CORAL, 0.85),
    ("Somatic Mutations", "frequency & type across top mutated genes", TEAL, 0.55),
    ("Pathway Activity", "GSVA oncogenic-pathway enrichment scores", BLUE, 0.65),
    ("Cell-State / Module / Clinical", "cell-type & gene-module enrichment + clinical traits", PURPLE, 0.7),
]
y0 = Inches(2.25)
bar_area_x = Inches(4.6)
max_bar_w = Inches(6.9)
for i, (t, d, c, frac) in enumerate(feat_rows):
    y = y0 + i * Inches(1.05)
    textbox(s, Inches(0.55), y, Inches(3.9), Inches(0.4), t, size=15, color=INK, bold=True)
    textbox(s, Inches(0.55), y + Inches(0.4), Inches(3.9), Inches(0.5), d, size=10.5, color=GRAY,
            line_spacing=1.1)
    rrect(s, bar_area_x, y + Inches(0.05), max_bar_w, Inches(0.42), GRAY_LT, radius=0.5)
    rrect(s, bar_area_x, y + Inches(0.05), Emu(int(max_bar_w * frac)), Inches(0.42), c, radius=0.5)

card(s, Inches(0.55), Inches(6.55), Inches(12.25), Inches(0.68))
textbox(s, Inches(0.8), Inches(6.68), Inches(11.8), Inches(0.45),
        "Spoiler for later: gene expression alone recovers 98.8% of full-model accuracy.",
        size=13, color=TEAL, bold=True, italic=True)
progress_dots(s, 2)
footer(s)

# ============================================================================
# SLIDE 10 — INGREDIENT 2: THE DRUG  (Act II)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 2, "Ingredient #2")
slide_title(s, "The drug: teaching the model to read a molecule")

picture_fit(s, os.path.join(MEDIA, "TFLSTM-1.png"), Inches(0.55), Inches(1.45), Inches(7.6), Inches(3.35))
textbox(s, Inches(0.55), Inches(4.85), Inches(7.6), Inches(0.85),
        "KD-Embed: a TF-LSTM autoencoder is pretrained on 2.5M compounds (96.7% valid SMILES"
        " reconstruction), then its latent code is transferred, zero-shot, onto every BeatAML drug.",
        size=12.5, color=GRAY, line_spacing=1.2)

reps = [
    ("PC", "473 physico-chemical\ndescriptors + fingerprints (RDKit)"),
    ("KD-Embed", "learned latent code\n(this study's best performer)"),
    ("ChemBERTa", "384-dim chemical\nlanguage model (77M compounds)"),
    ("MolFormer", "768-dim chemical\nlanguage model (1.1B SMILES)"),
]
rx = Inches(8.45); rw = Inches(4.35)
for i, (t, d) in enumerate(reps):
    y = Inches(1.45) + i * Inches(1.3)
    fill = GOLD if t == "KD-Embed" else CARD_BG
    txtcol = NAVY if t == "KD-Embed" else INK
    card(s, rx, y, rw, Inches(1.12), fill=fill)
    textbox(s, rx + Inches(0.2), y + Inches(0.12), rw - Inches(0.4), Inches(0.35), t, size=15,
            color=txtcol, bold=True)
    textbox(s, rx + Inches(0.2), y + Inches(0.48), rw - Inches(0.4), Inches(0.6), d, size=11,
            color=(NAVY2 if t == "KD-Embed" else GRAY), line_spacing=1.1)

progress_dots(s, 2)
footer(s)

# ============================================================================
# SLIDE 11 — THE SECRET SAUCE: DRUG-PATHWAY PROXIMITY  (Act II)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 2, "A patient-specific twist on drug features")
slide_title(s, "The secret sauce: how close is this drug to what's broken?")

picture_fit(s, os.path.join(MEDIA, "DrugPathwayProximity-1.png"),
            Inches(0.55), Inches(1.5), Inches(12.25), Inches(4.55))
textbox(s, Inches(0.55), Inches(6.25), Inches(12.25), Inches(0.75),
        "Drug targets propagate across the STRING protein network (random walk with restart); proximity"
        " to each patient's active oncogenic pathway becomes a bespoke drug-pathway-AUC feature.",
        size=13, color=GRAY, line_spacing=1.2)
progress_dots(s, 2)
footer(s)

# ============================================================================
# SLIDE 12 — THE BRAIN: WHY TABPFN  (Act II)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 2, "The model behind the curtain")
slide_title(s, "The brain: why a foundation model for tables?")

boxes = [
    ("Pretrained prior", "TabPFN is pre-trained once,\non millions of synthetic\ntabular tasks", CORAL),
    ("In-context learning", "Training patients become\ncontext — no gradient\ndescent, no per-task tuning", TEAL),
    ("One forward pass", "A new patient-drug pair\nis predicted instantly,\nadapting to the data at hand", GOLD),
]
bw = Inches(3.7); gapb = Inches(0.35); startxb = Inches(0.85); yb = Inches(1.85)
for i, (t, d, c) in enumerate(boxes):
    x = startxb + i * (bw + gapb)
    card(s, x, yb, bw, Inches(2.1))
    oval(s, x + Inches(0.25), yb + Inches(0.25), Inches(0.5), Inches(0.5), c)
    textbox(s, x + Inches(0.25), yb + Inches(0.95), bw - Inches(0.5), Inches(0.4), t, size=16,
            color=INK, bold=True)
    textbox(s, x + Inches(0.25), yb + Inches(1.35), bw - Inches(0.5), Inches(0.65), d, size=11.5,
            color=GRAY, line_spacing=1.15)
    if i < 2:
        connector_arrow(s, x + bw + Inches(0.02), yb + Inches(1.05), x + bw + Inches(0.33),
                         yb + Inches(1.05), color=GRAY, width=1.75)

textbox(s, Inches(0.85), Inches(4.35), Inches(11), Inches(0.4),
        "MEASURED AGAINST 10 OTHER MODEL FAMILIES", size=12.5, color=INK, bold=True)
zoo = ["Ridge", "LSSVR", "Random\nForest", "XGBoost", "LightGBM", "CatBoost", "CNN", "LSTM",
       "GAT", "GCN", "TabPFN"]
zw = Inches(1.02); zgap = Inches(0.06); zx = Inches(0.85); zy = Inches(4.85)
for i, name in enumerate(zoo):
    x = zx + i * (zw + zgap)
    fill = GOLD if name == "TabPFN" else GRAY_LT
    txtc = NAVY if name == "TabPFN" else INK
    chip(s, x, zy, zw, Inches(0.85), name, fill, txt_color=txtc, size=11)

textbox(s, Inches(0.85), Inches(5.95), Inches(11), Inches(0.7),
        "TabPFN wins not by the biggest cross-validation score, but by the smallest gap"
        " between cross-validation and truly held-out patients (next up).",
        size=13.5, color=BLUE, italic=True, bold=True, line_spacing=1.2)
progress_dots(s, 2)
footer(s)

print("slides 1-12 built")

# Each CV strategy is shown as a small patient(rows) x drug(columns)
# response-matrix grid, colored train=blue / validation=gold, so the reader
# can literally see how each scheme partitions the matrix:
#   Random CV        -> cells scattered independent of row/col (leakage)
#   Patient-Stratified-> whole ROWS (patients) assigned together
#   Drug-Stratified   -> whole COLUMNS (drugs) assigned together
TRAIN_C = BLUE
VAL_C = GOLD
GRID_ROWS, GRID_COLS = 6, 10
CELL = Inches(0.24)
CGAP = Inches(0.045)
GRID_W = GRID_COLS * CELL + (GRID_COLS - 1) * CGAP
GRID_H = GRID_ROWS * CELL + (GRID_ROWS - 1) * CGAP

def cv_matrix(slide, grid_x, grid_top, pattern, leak_rows=()):
    for r, row in enumerate(pattern):
        is_leak = r in leak_rows
        for c, val in enumerate(row):
            color = TRAIN_C if val == 1 else VAL_C
            x = grid_x + c * (CELL + CGAP)
            y = grid_top + r * (CELL + CGAP)
            rrect(slide, x, y, CELL, CELL, color, radius=0.3, line=(CORAL if is_leak else None))
    textbox(slide, grid_x, grid_top + GRID_H + Inches(0.02), GRID_W, Inches(0.18),
            "drugs →", size=8, color=GRAY, italic=True, align=PP_ALIGN.RIGHT)
    textbox(slide, grid_x, grid_top - Inches(0.19), GRID_W, Inches(0.16),
            "↓ patients", size=8, color=GRAY, italic=True, align=PP_ALIGN.LEFT)

def cv_legend(slide, x, y):
    sw = Inches(0.15)
    rrect(slide, x, y + Inches(0.02), sw, sw, TRAIN_C, radius=0.3)
    textbox(slide, x + sw + Inches(0.08), y, Inches(1.0), Inches(0.25), "train", size=10.5, color=INK)
    x2 = x + Inches(1.55)
    rrect(slide, x2, y + Inches(0.02), sw, sw, VAL_C, radius=0.3)
    textbox(slide, x2 + sw + Inches(0.08), y, Inches(1.2), Inches(0.25), "validation", size=10.5, color=INK)

# ============================================================================
# SLIDE 13 — THREE WAYS TO SPLIT THE DECK  (Act III)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 3, "Evaluation design decides everything")
slide_title(s, "How you split the data changes the whole verdict")

RANDOM_PATTERN = [
    [1,0,1,1,0,1,1,0,1,1],
    [0,1,1,0,1,0,1,1,0,1],
    [1,1,0,1,1,0,0,1,1,0],
    [0,1,1,1,0,1,1,0,1,1],
    [1,0,1,0,1,1,0,1,0,1],
    [1,1,0,1,1,0,1,1,0,0],
]
PS_PATTERN = [[1]*10, [1]*10, [1]*10, [1]*10, [0]*10, [0]*10]
DS_PATTERN = [[1,1,1,1,1,1,1,0,0,0]] * 6

cvs = [
    ("Random CV", "Patient-drug pairs shuffled freely.\nSame patient leaks into both\ntrain & validation.", CORAL,
     RANDOM_PATTERN, [1], "⚠ same patient, split both ways", CORAL, "looks great, means little"),
    ("Patient-Stratified CV", "Whole patients withheld.\nSimulates a brand-new\nperson walking in.", GOLD,
     PS_PATTERN, [], "✓ whole patient ROWS together", GREEN, "★ the clinical benchmark"),
    ("Drug-Stratified CV", "Whole drugs withheld.\nTests generalization to a\nstructurally novel compound.", TEAL,
     DS_PATTERN, [], "✓ whole drug COLUMNS together", GREEN, "the hard, honest frontier"),
]
cw = Inches(3.85); gapc = Inches(0.28); startxc = Inches(0.55); yc = Inches(1.75)
margin = Inches(0.25)
for i, (t, d, c, pattern, leak_rows, caption, cap_color, tag) in enumerate(cvs):
    x = startxc + i * (cw + gapc)
    card(s, x, yc, cw, Inches(4.55))
    rect(s, x, yc, cw, Inches(0.09), c)
    textbox(s, x + margin, yc + Inches(0.28), cw - 2*margin, Inches(0.45), t, size=17,
            color=INK, bold=True)
    textbox(s, x + margin, yc + Inches(0.8), cw - 2*margin, Inches(0.85), d, size=12,
            color=GRAY, line_spacing=1.2)
    textbox(s, x + margin, yc + Inches(1.63), cw - 2*margin, Inches(0.22), caption, size=10,
            color=cap_color, bold=True, align=PP_ALIGN.CENTER)
    grid_x = x + (cw - GRID_W) // 2
    grid_top = yc + Inches(1.88)
    cv_matrix(s, grid_x, grid_top, pattern, leak_rows=leak_rows)
    cv_legend(s, x + margin + Inches(0.35), yc + Inches(3.6))
    chip(s, x + margin, yc + Inches(4.03), cw - 2*margin, Inches(0.4), tag, c, size=11.5)

progress_dots(s, 3)
footer(s)

# ============================================================================
# SLIDE 14 — RANDOM CV LIES  (Act III)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 3, "The plot twist")
slide_title(s, "Random CV lies. Patient-Stratified CV tells the truth.")

cd2 = CategoryChartData()
cd2.categories = ["Random CV\n(inflated)", "Patient-Stratified CV\n(honest estimate)", "Held-out Test\n(the real world)"]
cd2.add_series("Pearson r", (0.807, 0.724, 0.690))
gx2, gy2, gw3, gh3 = Inches(0.75), Inches(1.75), Inches(7.6), Inches(4.7)
gf2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, gx2, gy2, gw3, gh3, cd2)
ch2 = gf2.chart
ch2.has_legend = False; ch2.has_title = False
pl2 = ch2.plots[0]; pl2.gap_width = 50
sr2 = pl2.series[0]
colorlist = [CORAL, GOLD, TEAL]
for i, pt in enumerate(sr2.points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = colorlist[i]
ca2 = ch2.category_axis; ca2.tick_labels.font.size = Pt(12.5); ca2.tick_labels.font.color.rgb = INK
va2 = ch2.value_axis; va2.minimum_scale = 0; va2.maximum_scale = 0.9
va2.has_major_gridlines = False
va2.tick_labels.font.size = Pt(11); va2.tick_labels.font.color.rgb = GRAY

card(s, Inches(8.65), Inches(1.75), Inches(4.15), Inches(4.7))
textbox(s, Inches(8.9), Inches(2.0), Inches(3.7), Inches(0.5), "WHY  THE  GAP?", size=13,
        color=INK, bold=True)
textbox(s, Inches(8.9), Inches(2.55), Inches(3.7), Inches(1.8),
        "Random CV lets different drug-responses from the SAME patient sit in both training"
        " and validation. The model quietly memorizes the patient instead of learning to"
        " generalize.",
        size=12.5, color=GRAY, line_spacing=1.25)
rect(s, Inches(8.9), Inches(4.5), Inches(3.4), Inches(0.02), GRAY_LT)
bignum(s, Inches(8.9), Inches(4.7), Inches(3.6), "≈15%", "relative inflation in Pearson r\nfrom Random CV vs. Test",
       CORAL, num_size=40, label_size=12)

textbox(s, Inches(0.75), Inches(6.6), Inches(11.9), Inches(0.5),
        "PS-CV tracks held-out test almost exactly (Δ ≈ 4.7%) — it is the benchmark we trust throughout this study.",
        size=13, color=BLUE, bold=True, italic=True)
progress_dots(s, 3)
footer(s)

# ============================================================================
# SLIDE 15 — THE SHOWDOWN: MODEL ZOO  (Act III)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 3, "Benchmarking under the honest CV")
slide_title(s, "The showdown: TabPFN vs. ten other model families")

picture_fit(s, os.path.join(MEDIA, "crop_fig4_panelA.png"), Inches(0.4), Inches(1.5), Inches(6.3), Inches(4.55))
picture_fit(s, os.path.join(MEDIA, "crop_fig4_panelB.png"), Inches(6.65), Inches(1.5), Inches(6.3), Inches(4.55))
textbox(s, Inches(0.4), Inches(1.35), Inches(6.3), Inches(0.35), "Classical ML vs. TabPFN", size=13,
        color=INK, bold=True, align=PP_ALIGN.CENTER)
textbox(s, Inches(6.65), Inches(1.35), Inches(6.3), Inches(0.35), "Deep Learning vs. TabPFN", size=13,
        color=INK, bold=True, align=PP_ALIGN.CENTER)

card(s, Inches(0.55), Inches(6.35), Inches(12.25), Inches(0.68))
textbox(s, Inches(0.8), Inches(6.48), Inches(11.8), Inches(0.45),
        "Circles = held-out test value. CNN/LSTM post higher CV bars but drop further at test — TabPFN's"
        " circle sits right where its bar says it should.",
        size=12.5, color=INK, italic=True)
progress_dots(s, 3)
footer(s)

# ============================================================================
# SLIDE 16 — THE WINNER: HEADLINE RESULT  (Act III)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 3, "The headline result")
slide_title(s, "TabPFN + KD-Embed: the winning combination")

picture_fit(s, os.path.join(MEDIA, "crop_shap_panelA.png"), Inches(0.55), Inches(1.55), Inches(6.7), Inches(4.9))
textbox(s, Inches(0.55), Inches(6.42), Inches(6.7), Inches(0.5),
        "Full held-out test set — 18,826 never-before-seen patient-drug pairs.",
        size=12, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

bx = Inches(7.7)
bignum(s, bx, Inches(1.7), Inches(4.9), "r = 0.690", "Pearson correlation, held-out test\n(vs. 0.36 ElasticNet · 0.68 MDREAM)",
       GOLD, num_size=46, label_size=13)
bignum(s, bx, Inches(3.35), Inches(4.9), "MAE = 37.7", "mean absolute error, AUC units\n(scale 0-800)", TEAL,
       num_size=46, label_size=13)
bignum(s, bx, Inches(5.0), Inches(4.9), "r = 0.166", "bonus: still significant on 14 drugs\n+ patients never seen at all (p<1e-5)",
       PURPLE, num_size=46, label_size=13)

progress_dots(s, 3)
footer(s)

# ============================================================================
# SLIDE 17 — WHAT MATTERS MOST: ABLATION  (Act III)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 3, "Feature ablation study")
slide_title(s, "What matters most? A reduced panel almost as good as all four")

cd3 = CategoryChartData()
cd3.categories = ["Full model\n(4 groups)", "Gene Expr.\nonly", "Mutation\nonly", "Pathway\nonly",
                  "Cell/Clin/Mod\nonly", "Mutation +\nCell/Clin/Mod"]
cd3.add_series("Test Pearson r", (0.690, 0.682, 0.647, 0.666, 0.665, 0.689))
gx3, gy3, gw4, gh4 = Inches(0.55), Inches(1.65), Inches(9.0), Inches(4.7)
gf3 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, gx3, gy3, gw4, gh4, cd3)
ch3 = gf3.chart
ch3.has_legend = False; ch3.has_title = False
pl3 = ch3.plots[0]; pl3.gap_width = 40
sr3 = pl3.series[0]
cols3 = [GOLD, CORAL, GRAY_LT, GRAY_LT, GRAY_LT, TEAL]
for i, pt in enumerate(sr3.points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = cols3[i]
ca3 = ch3.category_axis; ca3.tick_labels.font.size = Pt(11); ca3.tick_labels.font.color.rgb = INK
va3 = ch3.value_axis; va3.minimum_scale = 0.55; va3.maximum_scale = 0.72
va3.has_major_gridlines = False
va3.tick_labels.font.size = Pt(11); va3.tick_labels.font.color.rgb = GRAY

card(s, Inches(9.75), Inches(1.65), Inches(3.05), Inches(4.7))
textbox(s, Inches(9.95), Inches(1.85), Inches(2.7), Inches(0.7), "TWO  NOVEL  TAKEAWAYS", size=12,
        color=INK, bold=True, line_spacing=1.1)
textbox(s, Inches(9.95), Inches(2.55), Inches(2.7), Inches(1.6),
        "Gene expression alone recovers 98.8% of full-model accuracy.",
        size=12, color=CORAL, bold=True, line_spacing=1.2)
textbox(s, Inches(9.95), Inches(4.15), Inches(2.7), Inches(1.9),
        "Mutation + cell-state alone MATCHES the full model — a resource-light panel for"
        " clinics without deep RNA-seq.",
        size=12, color=TEAL, bold=True, line_spacing=1.2)

progress_dots(s, 3)
footer(s)

print("slides 1-17 built")

# ============================================================================
# SLIDE 18 — CRACKING THE BLACK BOX: SHAP  (Act IV)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 4, "Interpretability")
slide_title(s, "Cracking open the black box with SHAP")

picture_fit(s, os.path.join(MEDIA, "SHAP_beeswarm_test-1.png"), Inches(0.55), Inches(1.5), Inches(5.1), Inches(5.35))
textbox(s, Inches(0.55), Inches(6.85), Inches(5.1), Inches(0.35),
        "Top 30 features by mean |SHAP|, full test set", size=10.5, color=GRAY, italic=True,
        align=PP_ALIGN.CENTER)

rx = Inches(6.0)
bignum(s, rx, Inches(1.6), Inches(6.4), "23 / 30", "top features are KD-Embed latent\ndimensions — drug identity", CORAL,
       num_size=42, label_size=13)
bignum(s, rx, Inches(3.15), Inches(6.4), "7 / 30", "are patient genes: DNTT, CSPG4,\nCEACAM6, IGSF10, CDC42BPA...", TEAL,
       num_size=42, label_size=13)
card(s, rx, Inches(4.75), Inches(6.4), Inches(1.85))
textbox(s, rx + Inches(0.25), Inches(4.95), Inches(5.9), Inches(0.4), "THE  TWO-TIER  STORY", size=12.5,
        color=INK, bold=True)
textbox(s, rx + Inches(0.25), Inches(5.35), Inches(5.9), Inches(1.15),
        "Drug identity (KD-Embed) sets the population-level direction and magnitude of a"
        " prediction; a handful of patient genes then fine-tune it up or down for that individual.",
        size=13, color=GRAY, line_spacing=1.25)

progress_dots(s, 4)
footer(s)

# ============================================================================
# SLIDE 19 — READING A SINGLE PREDICTION  (Act IV)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 4, "From population to one person")
slide_title(s, "Reading a single prediction, feature by feature")

picture_fit(s, os.path.join(MEDIA, "crop_shap_panelC.png"), Inches(0.55), Inches(1.45), Inches(12.2), Inches(2.55))
picture_fit(s, os.path.join(MEDIA, "crop_shap_panelD.png"), Inches(0.55), Inches(4.15), Inches(12.2), Inches(2.55))

card(s, Inches(0.75), Inches(1.6), Inches(3.5), Inches(0.85), fill=RGBColor(0xFC,0xEC,0xE8))
textbox(s, Inches(0.95), Inches(1.72), Inches(3.15), Inches(0.65),
        "IDH2 mutation still pulls toward\nsensitivity (φ=-6.3) — matches known\nBCL-2 synthetic lethality biology.",
        size=10.5, color=CORAL, bold=True, line_spacing=1.15)

card(s, Inches(0.75), Inches(4.3), Inches(3.5), Inches(0.85), fill=RGBColor(0xE8,0xF5,0xF2))
textbox(s, Inches(0.95), Inches(4.42), Inches(3.15), Inches(0.65),
        "POU4F1 emerges as the strongest\npatient sensitivity signal — a novel,\ntestable AML biomarker hypothesis.",
        size=10.5, color=TEAL, bold=True, line_spacing=1.15)

progress_dots(s, 4)
footer(s)

# ============================================================================
# SLIDE 20 — PATIENT VS DRUG: WHO'S PREDICTABLE?  (Act IV)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 4, "Two ways to slice performance")
slide_title(s, "Some patients are easy. Some drugs are hard.")

picture_fit(s, os.path.join(MEDIA, "crop_perf_panelA.png"), Inches(0.4), Inches(1.5), Inches(6.3), Inches(3.55))
picture_fit(s, os.path.join(MEDIA, "crop_perf_panelB.png"), Inches(6.65), Inches(1.5), Inches(6.3), Inches(3.55))

lcard_w = Inches(6.0)
card(s, Inches(0.55), Inches(5.2), lcard_w, Inches(1.7))
bignum(s, Inches(0.8), Inches(5.35), Inches(5.5), "0.746", "mean per-patient r — 92.9% of patients\nexceed r > 0.5; 73.2% exceed r > 0.7",
       CORAL, num_size=34, label_size=12)
card(s, Inches(6.8), Inches(5.2), lcard_w, Inches(1.7))
bignum(s, Inches(7.05), Inches(5.35), Inches(5.5), "0.335", "mean per-drug r — only 15.4% of drugs\nexceed r > 0.5 (top: Venetoclax, 0.816)",
       TEAL, num_size=34, label_size=12)

textbox(s, Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.4),
        "Patient biology transfers broadly across drugs; drug-level accuracy depends more on each compound's own pharmacology.",
        size=11.5, color=GRAY, italic=True)
progress_dots(s, 4)
footer(s)

# ============================================================================
# SLIDE 21 — THE VENETOCLAX STORY: TOP-10 DRUG BIOLOGY  (Act IV)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 4, "Multi-omics drivers of the top-10 drugs")
slide_title(s, "The Venetoclax story: same biology, opposite verdicts")

textbox(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.0),
        "A “monocyte-like” transcriptional score predicts Venetoclax RESISTANCE (β=0.77, R²=0.87) — yet the same"
        " score predicts SENSITIVITY to MEK inhibitors and Panobinostat (β=-0.41 to -0.54). BCL6 expression and an"
        " immunogenic-cell-death (ICD) signature are significant across nearly all top-10 drugs — a new pan-drug axis.",
        size=13.5, color=INK, line_spacing=1.3)

picture_fit(s, os.path.join(MEDIA, "crop_perf_panelC.png"), Inches(0.55), Inches(2.65), Inches(12.2), Inches(4.4))
progress_dots(s, 4)
footer(s)

print("slides 1-21 built")

# ============================================================================
# SLIDE 22 — ROAD TEST: EXTERNAL VALIDATION  (Act V)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 5, "Beyond BeatAML, zero retraining")
slide_title(s, "Road test: two independent hospitals, one frozen model")

picture_fit(s, os.path.join(MEDIA, "LeeAML_best_config_scatter-1.png"), Inches(0.55), Inches(1.55), Inches(5.9), Inches(4.6))
picture_fit(s, os.path.join(MEDIA, "FIMMAML_Embed_scatter-1.png"), Inches(6.85), Inches(1.55), Inches(5.9), Inches(4.6))

chip(s, Inches(0.85), Inches(1.7), Inches(2.6), Inches(0.42), "LeeAML · USA · n=30", BLUE, size=11.5)
chip(s, Inches(7.15), Inches(1.7), Inches(3.0), Inches(0.42), "FIMM-AML · Finland · n=187", GOLD, size=11.5)

card(s, Inches(0.55), Inches(6.35), Inches(12.25), Inches(0.72))
textbox(s, Inches(0.8), Inches(6.46), Inches(11.8), Inches(0.55),
        "LeeAML: r = 0.597 (beats published MDREAM r = 0.59)   ·   FIMM-AML: |r| = 0.557 across a"
        " completely different assay (DSS vs. AUC) — same biology, still recognized.",
        size=13, color=INK, bold=True, italic=True)
progress_dots(s, 5)
footer(s)

# ============================================================================
# SLIDE 23 — WHO NEEDS PERSONALIZING?  (Act V)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 5, "Where personalization actually pays off")
slide_title(s, "Who actually needs a personalized choice?")

picture_fit(s, os.path.join(MEDIA, "Drug_Heterogeneity_3Cohorts_TabPFN-1.png"),
            Inches(0.55), Inches(1.45), Inches(12.25), Inches(3.85))

card(s, Inches(0.55), Inches(5.55), Inches(6.0), Inches(1.5), fill=RGBColor(0xFC,0xEC,0xE8))
textbox(s, Inches(0.8), Inches(5.7), Inches(5.5), Inches(0.4), "PRIORITIZE  FOR  SELECTION", size=12.5,
        color=CORAL, bold=True)
textbox(s, Inches(0.8), Inches(6.1), Inches(5.5), Inches(0.85),
        "Venetoclax & Panobinostat: highest patient-to-patient variability in every cohort —"
        " exactly where picking the right patient matters most.",
        size=12, color=INK, line_spacing=1.2)

card(s, Inches(6.8), Inches(5.55), Inches(6.0), Inches(1.5), fill=RGBColor(0xEF,0xF3,0xF8))
textbox(s, Inches(7.05), Inches(5.7), Inches(5.5), Inches(0.4), "SAFE  TO  DE-PRIORITIZE", size=12.5,
        color=BLUE, bold=True)
textbox(s, Inches(7.05), Inches(6.1), Inches(5.5), Inches(0.85),
        "Vismodegib & GDC-0879: near-uniform resistance for everyone (SD < 3) —"
        " model-guided selection adds nothing here.",
        size=12, color=INK, line_spacing=1.2)
progress_dots(s, 5)
footer(s)

# ============================================================================
# SLIDE 24 — HOW WE STACK UP + LIMITATIONS  (Act V)
# ============================================================================
s = add_slide()
bg(s)
kicker_and_bar(s, 5, "Positioning & honesty")
slide_title(s, "How PREDICT-AML stacks up — and where it still struggles")

textbox(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(0.4), "vs.  NetAML  (network-driven, 2025)",
        size=13.5, color=INK, bold=True)
advs = [
    "One unified model — not 87 separate per-drug models",
    "Four multi-omics layers, not RNA-seq alone",
    "Uses drug molecular structure; NetAML uses none",
    "Evaluated on a temporally-held-out cohort, not random 70/30",
    "Validated on 2 external cohorts (|r|=0.557) vs. NetAML's 1 (r=0.38)",
]
yA = Inches(2.0)
for t in advs:
    oval(s, Inches(0.55), yA + Inches(0.06), Inches(0.14), Inches(0.14), TEAL)
    textbox(s, Inches(0.85), yA, Inches(5.75), Inches(0.6), t, size=12.5, color=INK, line_spacing=1.15)
    yA += Inches(0.72)

rect(s, Inches(6.75), Inches(1.5), Inches(0.02), Inches(5.1), GRAY_LT)

textbox(s, Inches(7.1), Inches(1.5), Inches(5.7), Inches(0.4), "LIMITATIONS  WE  OWN", size=13.5,
        color=INK, bold=True)
lims = [
    "Novel-drug generalization remains hard (DS-CV r ≈ 0.37)",
    "FIMM-AML cross-platform check (AUC vs. DSS) is inherently noisy",
    "Predicts a continuous AUC, not yet a binary clinical decision",
    "POU4F1 and related hits are computational hypotheses,\nawaiting wet-lab validation",
]
yB = Inches(2.0)
for t in lims:
    oval(s, Inches(7.1), yB + Inches(0.06), Inches(0.14), Inches(0.14), CORAL)
    textbox(s, Inches(7.4), yB, Inches(5.4), Inches(0.75), t, size=12.5, color=INK, line_spacing=1.15)
    yB += Inches(0.85)

progress_dots(s, 5)
footer(s)

# ============================================================================
# SLIDE 25 — TAKE-HOME / CLOSING
# ============================================================================
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, SLIDE_W, Inches(0.09), GOLD_TXT)
textbox(s, Inches(0.7), Inches(0.55), Inches(8), Inches(0.4), "TAKE-HOME  MESSAGE", size=13,
        color=GOLD_TXT, bold=True)
textbox(s, Inches(0.65), Inches(1.05), Inches(11.8), Inches(1.1),
        "Precision oncology, one patient at a time.", size=38, color=WHITE, bold=True)

nums = [("0.690", "test Pearson r,\nnever-seen patients"), ("2", "independent external\ncohorts validated"),
        ("150", "drugs covered by\na single model"), ("~1,300", "multi-omic features,\ninterpretable via SHAP")]
nw = Inches(2.85); gapn = Inches(0.22); startxn = Inches(0.7); yn = Inches(2.55)
for i, (num, lab) in enumerate(nums):
    x = startxn + i * (nw + gapn)
    rrect(s, x, yn, nw, Inches(1.85), NAVY2, radius=0.08)
    textbox(s, x + Inches(0.2), yn + Inches(0.22), nw - Inches(0.4), Inches(0.7), num, size=32,
            color=GOLD_TXT, bold=True)
    textbox(s, x + Inches(0.2), yn + Inches(0.95), nw - Inches(0.4), Inches(0.75), lab, size=12,
            color=RGBColor(0xC7,0xD1,0xE6), line_spacing=1.15)

textbox(s, Inches(0.7), Inches(4.85), Inches(11.6), Inches(0.9),
        "Honest, patient-stratified evaluation. A model that knows the drug and knows the patient.\n"
        "Interpretable enough to teach us AML biology we didn't already know.",
        size=15.5, color=WHITE, line_spacing=1.3)

rect(s, Inches(0.72), Inches(6.0), Inches(0.55), Inches(0.04), GOLD_TXT)
textbox(s, Inches(0.7), Inches(6.2), Inches(9), Inches(0.4),
        "Code & workflow: github.com/raghvendra5686/PT-AML2.0", size=13, color=GOLD_TXT, bold=True)
textbox(s, Inches(0.7), Inches(6.6), Inches(9), Inches(0.4),
        "Al-Ani · Jani · Bensmail · Mall   —   Qatar Computing Research Institute, HBKU",
        size=12, color=GRAY_LT)
textbox(s, Inches(0.7), Inches(6.95), Inches(9), Inches(0.4), "Thank you.", size=14, color=WHITE,
        italic=True, bold=True)

footer(s, dark=True)

print("slides 1-25 built — DECK COMPLETE")

prs.save(OUT)
print("SAVED:", OUT)
print("Total slides:", len(prs.slides))

